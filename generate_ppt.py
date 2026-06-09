#!/usr/bin/env python3
"""
Generate a comprehensive PowerPoint presentation for VlogHub project
with project details and images
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Define color scheme
COLOR_PRIMARY = RGBColor(31, 78, 121)      # Dark blue
COLOR_ACCENT = RGBColor(192, 0, 0)         # Red
COLOR_TEXT = RGBColor(51, 51, 51)          # Dark gray
COLOR_LIGHT = RGBColor(242, 242, 242)      # Light gray

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(60)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = subtitle
    subtitle_p.font.size = Pt(28)
    subtitle_p.font.color.rgb = RGBColor(220, 220, 220)
    
    # Date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(0.5))
    date_frame = date_box.text_frame
    date_p = date_frame.paragraphs[0]
    date_p.text = "June 9, 2026"
    date_p.font.size = Pt(18)
    date_p.font.color.rgb = RGBColor(200, 200, 200)

def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Header bar
    header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = COLOR_PRIMARY
    header_shape.line.color.rgb = COLOR_PRIMARY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.3))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        p.level = 0

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Add a two-column content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Header bar
    header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = COLOR_PRIMARY
    header_shape.line.color.rgb = COLOR_PRIMARY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.0), Inches(4.5), Inches(5.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    # Left title
    p = left_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT
    
    for item in left_items:
        p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXT
        p.level = 0
        p.space_before = Pt(4)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.0), Inches(4.5), Inches(5.5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    # Right title
    p = right_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT
    
    for item in right_items:
        p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXT
        p.level = 0
        p.space_before = Pt(4)

def add_image_slide(prs, title, image_paths):
    """Add a slide with multiple images"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Header bar
    header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.7))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = COLOR_PRIMARY
    header_shape.line.color.rgb = COLOR_PRIMARY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.55))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Add images in grid
    img_count = len(image_paths)
    if img_count == 0:
        return
    
    # Calculate layout
    if img_count <= 4:
        cols = 2
        rows = 2
        img_width = 4.0
        img_height = 3.0
    else:
        cols = 3
        rows = 2
        img_width = 3.0
        img_height = 2.3
    
    x_start = 0.3
    y_start = 0.9
    x_spacing = 3.3 if cols == 3 else 4.8
    y_spacing = 2.5 if rows == 2 else 2.8
    
    for idx, img_path in enumerate(image_paths):
        if idx >= cols * rows:
            break
        if os.path.exists(img_path):
            col = idx % cols
            row = idx // cols
            x = x_start + col * x_spacing
            y = y_start + row * y_spacing
            try:
                slide.shapes.add_picture(img_path, Inches(x), Inches(y), width=Inches(img_width))
            except Exception as e:
                print(f"Could not add image {img_path}: {e}")

def create_vloghub_presentation():
    """Create the complete VlogHub presentation"""
    
    # Initialize presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    image_dir = r"next-vlogging\public\assets\images"
    image_files = [
        os.path.join(image_dir, f"thumb{i}.jpg") for i in range(1, 13)
    ]
    
    # Slide 1: Title
    add_title_slide(prs, "VlogHub", "A Complete Vlogging Platform")
    
    # Slide 2: Executive Summary
    add_content_slide(prs, "Executive Summary", [
        "✓ Full-stack web application for vloggers",
        "✓ Frontend: Next.js 16.2.1 with React 19.2.4",
        "✓ Backend: Django 6.0.4 with REST API",
        "✓ 20+ interactive pages & routes",
        "✓ OAuth integration (Google & Apple)",
        "✓ Real-time analytics & tracking system",
        "✓ E-commerce merch store (ready)",
        "✓ Live streaming with Super Chat"
    ])
    
    # Slide 3: Architecture Overview
    add_content_slide(prs, "Architecture", [
        "• Frontend: Next.js App Router with TypeScript",
        "• Backend: Django REST Framework API",
        "• Database: SQLite (dev) / PostgreSQL (prod)",
        "• State Management: Zustand stores",
        "• Styling: Tailwind CSS 4 + Framer Motion",
        "• Authentication: JWT + OAuth (Google/Apple)",
        "• Deployment: Vercel (frontend) + Render (backend)",
        "• Real-time: WebSocket ready for live features"
    ])
    
    # Slide 4: Tech Stack
    add_two_column_slide(prs, "Technology Stack",
        "Frontend", [
            "• Next.js 16.2.1",
            "• React 19.2.4",
            "• TypeScript 5",
            "• Tailwind CSS 4",
            "• Framer Motion",
            "• Zustand",
            "• Radix UI"
        ],
        "Backend", [
            "• Django 6.0.4",
            "• Python 3.14",
            "• Django REST",
            "• PostgreSQL",
            "• JWT Auth",
            "• Google OAuth 2.0",
            "• Apple Sign-In"
        ])
    
    # Slide 5: Database Models
    add_content_slide(prs, "Database Schema", [
        "🔹 User → Extended profile with avatars",
        "🔹 Video → Full metadata (title, description, URL)",
        "🔹 Comment → Nested comment threads",
        "🔹 TrackingEvent → User behavior analytics",
        "🔹 UserVideoState → Like/Dislike interactions",
        "🔹 PasswordResetToken → Secure password recovery",
        "Total: 6 models with optimized indexes",
        "Status: Ready for production PostgreSQL"
    ])
    
    # Slide 6: Frontend Pages Part 1
    add_two_column_slide(prs, "Frontend Pages (Phase 1)",
        "Core Content", [
            "✓ Home Page",
            "✓ Video Detail Player",
            "✓ Category View",
            "✓ About Page",
            "✓ Contact Page"
        ],
        "Engagement", [
            "✓ Shorts Feed",
            "✓ Community/Blog",
            "✓ Live Streaming",
            "✓ Creator Shop",
            "✓ Search Results"
        ])
    
    # Slide 7: Frontend Pages Part 2
    add_two_column_slide(prs, "Frontend Pages (Phase 2)",
        "User Account", [
            "✓ Login / Register",
            "✓ Password Recovery",
            "✓ User Profile",
            "✓ Watch Later Library",
            "✓ Series/Playlists"
        ],
        "Legal & Other", [
            "✓ Privacy Policy",
            "✓ Terms of Service",
            "✓ API Proxy Routes",
            "✓ Backend Testing",
            "✓ Error Handling"
        ])
    
    # Slide 8: Key Features
    add_content_slide(prs, "Key Features", [
        "🎥 Video Management: Upload, metadata, categorization",
        "🔐 Security: Email/OAuth authentication + JWT",
        "💬 Social: Comments, likes, shares, watch later",
        "📊 Analytics: Watch progress, interactions, views",
        "🛍️ Monetization: Merch shop, live super chat",
        "🎨 UX: Dark/light themes, mobile responsive",
        "⚡ Performance: Turbopack, database indexing",
        "🌐 Scalability: Production-ready architecture"
    ])
    
    # Slide 9: API Endpoints
    add_two_column_slide(prs, "API Endpoints",
        "Authentication", [
            "• POST /auth/register",
            "• POST /auth/login",
            "• POST /auth/google",
            "• POST /auth/apple",
            "• POST /auth/forgot-password",
            "• POST /auth/reset-password"
        ],
        "Content & Tracking", [
            "• GET/POST /videos",
            "• GET/POST /videos/[id]/comments",
            "• GET /videos/trending",
            "• GET /categories",
            "• POST /tracking/event",
            "• GET /tracking/stats"
        ])
    
    # Slide 10: Development Status
    add_two_column_slide(prs, "Development Status",
        "✅ Completed", [
            "✓ Frontend structure",
            "✓ Backend API",
            "✓ Database models",
            "✓ OAuth integration",
            "✓ Comment system",
            "✓ Theme support",
            "✓ Analytics tracking"
        ],
        "⚠️ Planned / 🔮 Future", [
            "⚠ Shop activation",
            "⚠ Live streaming",
            "⚠ Super Chat",
            "🔮 Admin dashboard",
            "🔮 Video upload",
            "🔮 Recommendations",
            "🔮 CDN integration"
        ])
    
    # Slide 11: Deployment
    add_two_column_slide(prs, "Deployment Architecture",
        "Production Stack", [
            "Frontend: Vercel",
            "Backend: Render",
            "Database: PostgreSQL",
            "  (Render/Railway/Neon)",
            "Static Files: WhiteNoise",
            "WSGI: Gunicorn",
            "CDN: Ready"
        ],
        "Environment Setup", [
            "• SECRET_KEY setup",
            "• CORS configuration",
            "• OAuth credentials",
            "• Database URL",
            "• Email service config",
            "• Domain settings",
            "• SSL/TLS enabled"
        ])
    
    # Slide 12: Project Statistics
    add_content_slide(prs, "Project Statistics", [
        "📊 Frontend Pages: 20+ routes",
        "📊 React Components: 25+ reusable",
        "📊 API Endpoints: 30+ endpoints",
        "📊 Database Models: 6 models",
        "📊 Authentication Methods: 3 (Email/Google/Apple)",
        "📊 NPM Dependencies: 15+",
        "📊 Python Dependencies: 8",
        "📊 Code Lines: 5,000+ (Frontend + Backend)"
    ])
    
    # Slide 13: Running the Project
    add_content_slide(prs, "Quick Start Guide", [
        "🚀 Backend: cd django-backend && python manage.py runserver",
        "🚀 Frontend: cd next-vlogging && npm run dev",
        "🌐 Frontend URL: http://localhost:3000",
        "🌐 Backend URL: http://127.0.0.1:8000",
        "🌐 Django Admin: http://127.0.0.1:8000/admin",
        "🌐 API Root: http://127.0.0.1:8000/api/hello/",
        "✓ Current Status: Both servers running",
        "✓ Ready for development & testing"
    ])
    
    # Slide 14-15: Project Images (Gallery)
    if len(image_files) >= 4:
        add_image_slide(prs, "Project Gallery - Set 1", image_files[:4])
    
    if len(image_files) >= 8:
        add_image_slide(prs, "Project Gallery - Set 2", image_files[4:8])
    
    if len(image_files) >= 12:
        add_image_slide(prs, "Project Gallery - Set 3", image_files[8:12])
    
    # Slide 16: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = "VlogHub"
    title_p.font.size = Pt(66)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = "A Modern Vlogging Platform\nBuilt with Next.js & Django"
    subtitle_p.font.size = Pt(28)
    subtitle_p.font.color.rgb = RGBColor(220, 220, 220)
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # Status
    status_box = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(8), Inches(1))
    status_frame = status_box.text_frame
    status_p = status_frame.paragraphs[0]
    status_p.text = "✓ Status: Running | 🚀 Active Development"
    status_p.font.size = Pt(20)
    status_p.font.color.rgb = RGBColor(100, 200, 100)
    status_p.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    output_path = "VlogHub_Project_Report.pptx"
    prs.save(output_path)
    print(f"✓ PowerPoint presentation created: {output_path}")
    return output_path

if __name__ == "__main__":
    create_vloghub_presentation()
